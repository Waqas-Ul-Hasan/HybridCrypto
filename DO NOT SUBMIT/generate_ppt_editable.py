import asyncio
import os
import re
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def parse_rgb(color_str):
    match = re.search(r'rgba?\((\d+),\s*(\d+),\s*(\d+)', color_str)
    if match:
        return RGBColor(int(match.group(1)), int(match.group(2)), int(match.group(3)))
    return RGBColor(255, 255, 255)

async def generate_ppt_editable():
    html_path = "file:///C:/FA23-BCS-A/IS/IS-Terminal/slides.html"
    out_pptx = "slides_fully_editable.pptx"
    
    print("Launching headless browser for Multi-Layer PPT Extraction...")
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        viewport_w = 1280
        viewport_h = 720
        page = await browser.new_page(viewport={'width': viewport_w, 'height': viewport_h})
        await page.goto(html_path)
        
        # Hide navigation UI
        await page.evaluate("""
            document.querySelector('.nav').style.display = 'none';
            document.querySelector('.slide-counter').style.display = 'none';
        """)
        
        total = await page.evaluate("total")
        prs = Presentation()
        prs.slide_width = Inches(13.333333333)
        prs.slide_height = Inches(7.5)
        
        scale_x = prs.slide_width / viewport_w
        scale_y = prs.slide_height / viewport_h
        
        blank_slide_layout = prs.slide_layouts[6]
        
        for i in range(total):
            await page.evaluate(f"goTo({i})")
            await page.wait_for_timeout(600)
            
            print(f"  -> Processing Slide {i+1}/{total} layers...")
            
            # 1. Tag graphical elements
            graphical_js = """
            () => {
                let elements = [];
                document.querySelectorAll('.card, .tag, .code-block, .flow-box, .glow-line, .pbar-wrap, table, .ico').forEach((el, idx) => {
                    const slide = el.closest('.slide');
                    if (!slide || !slide.classList.contains('active')) return;
                    
                    const rect = el.getBoundingClientRect();
                    if (rect.width === 0 || rect.height === 0) return;
                    
                    let elId = 'export-gfx-' + idx;
                    el.setAttribute('data-export-gfx', elId);
                    elements.push({
                        id: elId,
                        x: rect.x,
                        y: rect.y,
                        w: rect.width,
                        h: rect.height
                    });
                });
                return elements;
            }
            """
            gfx_data = await page.evaluate(graphical_js)
            
            # 2. Hide ALL graphics AND text to take pure background screenshot
            await page.evaluate("""
            () => {
                document.querySelectorAll('[data-export-gfx]').forEach(el => el.style.setProperty('visibility', 'hidden', 'important'));
                document.querySelectorAll('h1, h2, h3, p, li, th, td, .s1-lock, .big-check, .flow-arrow').forEach(el => {
                    const slide = el.closest('.slide');
                    if (!slide || !slide.classList.contains('active')) return;
                    el.style.setProperty('visibility', 'hidden', 'important');
                });
            }
            """)
            
            bg_path = f"bg_{i}.png"
            await page.screenshot(path=bg_path)
            
            # 3. Restore visibility, set page to transparent, hide text
            text_js = """
            () => {
                // Restore visibility
                document.querySelectorAll('*').forEach(el => el.style.removeProperty('visibility'));
                
                // Set page backgrounds to transparent
                document.body.style.setProperty('background', 'transparent', 'important');
                document.documentElement.style.setProperty('background', 'transparent', 'important');
                document.querySelectorAll('.slide').forEach(s => s.style.setProperty('background', 'transparent', 'important'));
                
                // Collect text AND hide it
                const text_elements = [];
                const selectors = 'h1, h2, h3, p, li, th, td, .tag, .s1-lock, .big-check, .flow-box, .sec-label > span, .code-block, .flow-arrow, .ico';
                
                document.querySelectorAll(selectors).forEach(el => {
                    const slide = el.closest('.slide');
                    if (!slide || !slide.classList.contains('active')) return;
                    
                    const rect = el.getBoundingClientRect();
                    if (rect.width === 0 || rect.height === 0) return;
                    
                    const style = window.getComputedStyle(el);
                    // Extract node text safely
                    let text = Array.from(el.childNodes)
                        .filter(node => node.nodeType === Node.TEXT_NODE)
                        .map(node => node.textContent)
                        .join('').trim();
                        
                    // fallback to innerText if no direct text nodes (e.g. nested spans)
                    if (!text) text = el.innerText.trim();
                    if (!text) return;
                    
                    text_elements.push({
                        text: text,
                        x: rect.x,
                        y: rect.y,
                        w: rect.width,
                        h: rect.height,
                        fontSize: style.fontSize,
                        color: style.color,
                        textAlign: style.textAlign,
                        fontWeight: style.fontWeight
                    });
                    
                    el.style.setProperty('color', 'transparent', 'important');
                    el.style.setProperty('text-shadow', 'none', 'important');
                    el.querySelectorAll('*').forEach(child => {
                        child.style.setProperty('color', 'transparent', 'important');
                        child.style.setProperty('text-shadow', 'none', 'important');
                    });
                });
                return text_elements;
            }
            """
            text_data = await page.evaluate(text_js)
            
            # 4. Screenshot each graphical element independently
            gfx_images = []
            for gfx in gfx_data:
                gfx_path = f"gfx_{i}_{gfx['id']}.png"
                # Hide child graphics to prevent baked-in layers
                await page.evaluate(f"""
                () => {{
                    const el = document.querySelector("[data-export-gfx='{gfx['id']}']");
                    if (el) {{
                        el.querySelectorAll('[data-export-gfx]').forEach(child => {{
                            child.setAttribute('data-hidden-temp', 'true');
                            child.style.setProperty('visibility', 'hidden', 'important');
                        }});
                    }}
                }}
                """)
                
                await page.locator(f"[data-export-gfx='{gfx['id']}']").screenshot(path=gfx_path, omit_background=True)
                
                # Unhide child graphics
                await page.evaluate(f"""
                () => {{
                    const el = document.querySelector("[data-export-gfx='{gfx['id']}']");
                    if (el) {{
                        el.querySelectorAll('[data-hidden-temp="true"]').forEach(child => {{
                            child.removeAttribute('data-hidden-temp');
                            child.style.removeProperty('visibility');
                        }});
                    }}
                }}
                """)
                
                gfx['path'] = gfx_path
                gfx_images.append(gfx)
            
            # 5. Revert styles for next slide
            await page.evaluate("""
            () => {
                document.querySelectorAll('*').forEach(el => {
                    el.style.removeProperty('color');
                    el.style.removeProperty('text-shadow');
                });
                document.body.style.removeProperty('background');
                document.documentElement.style.removeProperty('background');
                document.querySelectorAll('.slide').forEach(s => s.style.removeProperty('background'));
            }
            """)
            
            # 6. Construct PPT Slide
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Layer 1: Background
            slide.shapes.add_picture(bg_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            os.remove(bg_path)
            
            # Layer 2: Independent Graphics
            for gfx in gfx_images:
                try:
                    slide.shapes.add_picture(
                        gfx['path'], 
                        Inches(gfx['x'] / viewport_w * 13.333333333), 
                        Inches(gfx['y'] / viewport_h * 7.5),
                        width=Inches(gfx['w'] / viewport_w * 13.333333333),
                        height=Inches(gfx['h'] / viewport_h * 7.5)
                    )
                except Exception as e:
                    print(f"Warning: Could not add graphic {gfx['path']}: {e}")
                if os.path.exists(gfx['path']):
                    os.remove(gfx['path'])
                    
            # Layer 3: Text
            for item in text_data:
                x = int(item['x'] * scale_x)
                y = int(item['y'] * scale_y)
                w = int(item['w'] * scale_x)
                h = int((item['h'] + 10) * scale_y)
                
                txBox = slide.shapes.add_textbox(x, y, w, h)
                tf = txBox.text_frame
                tf.word_wrap = True
                tf.margin_left = 0
                tf.margin_top = 0
                tf.margin_right = 0
                tf.margin_bottom = 0
                
                p = tf.paragraphs[0]
                p.text = item['text']
                font_size_px = float(item['fontSize'].replace('px', ''))
                p.font.size = Pt(font_size_px * 0.75)
                p.font.color.rgb = parse_rgb(item['color'])
                
                if item['fontWeight'] in ['bold', '600', '700', '800', '900']:
                    p.font.bold = True
                    
                align = item['textAlign']
                if align == 'center':
                    p.alignment = PP_ALIGN.CENTER
                elif align == 'right':
                    p.alignment = PP_ALIGN.RIGHT
                else:
                    p.alignment = PP_ALIGN.LEFT
                    
        await browser.close()
        
    prs.save(out_pptx)
    print(f"\n[OK] Fully Modular Editable PPT saved -> {out_pptx}")

if __name__ == "__main__":
    asyncio.run(generate_ppt_editable())
