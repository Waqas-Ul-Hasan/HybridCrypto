import asyncio
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches
import os

async def generate_ppt_visual():
    # Use absolute path to ensure Playwright loads it correctly
    html_path = "file:///C:/FA23-BCS-A/IS/IS-Terminal/slides.html"
    out_pptx = "slides_v2.pptx"
    
    print("Launching headless browser to capture slides...")
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        # 16:9 format, but smaller resolution so the CSS text scaling appears larger
        page = await browser.new_page(viewport={'width': 1280, 'height': 720})
        await page.goto(html_path)
        
        # Hide the bottom navigation buttons and dots so they don't appear in the PPT
        await page.evaluate("""
            document.querySelector('.nav').style.display = 'none';
            document.querySelector('.slide-counter').style.display = 'none';
        """)
        
        # Grab the total number of slides from the global JS variable
        total = await page.evaluate("total")
        
        prs = Presentation()
        # Set presentation size to 16:9
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        blank_slide_layout = prs.slide_layouts[6] # Blank layout
        
        print(f"Found {total} slides to convert.")
        
        for i in range(total):
            # Tell the JS engine to switch to slide i
            await page.evaluate(f"goTo({i})")
            # Wait for CSS transitions to settle (450ms transition in css, wait 600ms)
            await page.wait_for_timeout(600)
            
            # Take screenshot of the viewport
            img_path = f"temp_slide_{i}.png"
            await page.screenshot(path=img_path)
            
            # Append a blank slide and stretch the screenshot over it
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            # Clean up the temp image
            os.remove(img_path)
            print(f"  -> Captured slide {i+1}/{total}")
            
        await browser.close()
        
    prs.save(out_pptx)
    print(f"\n[OK] High-fidelity visual presentation saved -> {out_pptx}")

if __name__ == "__main__":
    asyncio.run(generate_ppt_visual())
