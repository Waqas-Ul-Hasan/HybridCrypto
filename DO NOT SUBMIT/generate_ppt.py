from pptx import Presentation
from bs4 import BeautifulSoup
import os

def html_to_ppt(html_path, ppt_path):
    with open(html_path, 'r', encoding='utf-8') as f:
        html = f.read()
    
    soup = BeautifulSoup(html, 'html.parser')
    prs = Presentation()
    
    slides = soup.find_all('div', class_='slide')
    
    for slide_div in slides:
        is_title_slide = 'slide-1' in slide_div.get('class', [])
        
        if is_title_slide:
            slide_layout = prs.slide_layouts[0] # Title slide
        else:
            slide_layout = prs.slide_layouts[1] # Title and content
            
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_elem = slide_div.find(['h1', 'h2'])
        if title_elem and slide.shapes.title:
            # Remove any hidden text or script tags
            slide.shapes.title.text = title_elem.get_text(separator=' ', strip=True)
            
        # Content
        if len(slide.shapes.placeholders) > 1:
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            
            # Find elements that should be bullet points
            # We will grab h3, p, and li
            seen = set()
            for el in slide_div.find_all(['h3', 'p', 'li', 'td', 'div']):
                if el == title_elem:
                    continue
                # For divs, only capture code-block or specific things, else skip
                if el.name == 'div' and 'code-block' not in el.get('class', []):
                    continue
                
                # If we are inside a table, maybe extract td
                if el.name == 'p' and el.find_parent('td'):
                    continue
                    
                text = el.get_text(separator=' ', strip=True)
                if not text or text in seen:
                    continue
                seen.add(text)
                
                p = tf.add_paragraph()
                p.text = text
                if el.name == 'h3' or el.name == 'div':
                    p.level = 0
                else:
                    p.level = 1
                    
    prs.save(ppt_path)
    print(f"Saved {ppt_path}")

if __name__ == '__main__':
    html_to_ppt('slides.html', 'slides.pptx')
